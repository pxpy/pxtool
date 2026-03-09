from pptx import Presentation

def extract_ppt_text(file_path):
    prs = Presentation(file_path)
    content = []
    for i, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {i+1} ---"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += "\n" + shape.text
        content.append(slide_text)
    content = "\n".join(content)
    return len(content), content

# 将提取到的字符串放入 API 的 'content' 字段发送
if __name__ == "__main__":
    count, content  = extract_ppt_text('path')
    with open("txt/ppt内容.txt", "w", encoding="utf-8") as f:
        f.write(content)